"""
Smart Email Organizer — Professional PPT Generator
Creates a visually stunning presentation with custom design elements.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───
DARK_BG       = RGBColor(0x0F, 0x0F, 0x1A)   # Deep dark navy
CARD_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # Dark card
ACCENT_BLUE   = RGBColor(0x66, 0x7E, 0xEA)   # Primary accent
ACCENT_PURPLE = RGBColor(0x76, 0x4B, 0xA2)   # Secondary accent
ACCENT_PINK   = RGBColor(0xF0, 0x93, 0xFB)   # Highlight pink
ACCENT_CYAN   = RGBColor(0x4E, 0xCA, 0xDC)   # Cyan
ACCENT_GREEN  = RGBColor(0x48, 0xBB, 0x78)   # Green
ACCENT_ORANGE = RGBColor(0xED, 0x8A, 0x36)   # Orange
ACCENT_RED    = RGBColor(0xFC, 0x56, 0x56)   # Red
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB0, 0xC0)
SEMI_WHITE    = RGBColor(0xE0, 0xE0, 0xF0)
GRADIENT_START = RGBColor(0x66, 0x7E, 0xEA)
GRADIENT_END   = RGBColor(0x76, 0x4B, 0xA2)

# Category colors
CAT_EVENTS     = RGBColor(0xED, 0x8A, 0x36)
CAT_ACADEMICS  = RGBColor(0x48, 0xBB, 0x78)
CAT_HACKATHONS = RGBColor(0xF0, 0x93, 0xFB)
CAT_PERSONAL   = RGBColor(0x4E, 0xCA, 0xDC)
CAT_SPAM       = RGBColor(0xFC, 0x56, 0x56)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helper Functions ───

def add_dark_bg(slide):
    """Add dark gradient-like background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_gradient_bar(slide, left=0, top=0, width=None, height=Inches(0.06), color1=GRADIENT_START, color2=GRADIENT_END):
    """Add a gradient accent bar at the top of the slide."""
    if width is None:
        width = prs.slide_width
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()

def add_accent_shape(slide, left, top, width, height, color, alpha=0.15):
    """Add a decorative accent circle/shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.fore_color.brightness = 0.6
    shape.line.fill.background()

def add_card(slide, left, top, width, height, color=CARD_BG, border_color=None):
    """Add a card-like rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE, spacing=Pt(8)):
    """Add a bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_icon_circle(slide, left, top, size, color, emoji_text=""):
    """Add a colored circle with text/emoji."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if emoji_text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = emoji_text
        p.font.size = Pt(int(size / Inches(1) * 16))
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_slide_number(slide, num, total):
    """Add slide number indicator."""
    add_text_box(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)

# Top gradient bar
add_gradient_bar(slide)

# Large decorative circles (background)
add_accent_shape(slide, Inches(-2), Inches(-2), Inches(6), Inches(6), ACCENT_BLUE)
add_accent_shape(slide, Inches(10), Inches(4), Inches(5), Inches(5), ACCENT_PURPLE)
add_accent_shape(slide, Inches(5), Inches(5), Inches(4), Inches(4), ACCENT_PINK)

# Central content card
card = add_card(slide, Inches(2.5), Inches(1.5), Inches(8.3), Inches(4.5), 
                color=RGBColor(0x15, 0x15, 0x28), border_color=ACCENT_BLUE)

# Email icon
icon = add_icon_circle(slide, Inches(5.9), Inches(1.8), Inches(1.2), ACCENT_BLUE, "✉")

# Title
add_text_box(slide, Inches(2.5), Inches(3.2), Inches(8.3), Inches(1.0),
             "Smart Email Organizer", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

# Subtitle
add_text_box(slide, Inches(3.0), Inches(4.2), Inches(7.3), Inches(0.6),
             "AI-Powered Email Management for University Students", 
             font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(3.5), Inches(4.9), Inches(6.3), Inches(0.5),
             "Hybrid ML Classification  •  AI Chatbot Assistant  •  Smart Calendar",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom info
add_text_box(slide, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.6),
             "Information Systems Project  |  2026",
             font_size=16, color=SEMI_WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(11), Inches(-1), Inches(4), Inches(4), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "The Problem", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Why university students need smarter email management",
             font_size=16, color=ACCENT_BLUE)

# Problem cards
problems = [
    ("📧", "Email Overload", "Students receive 50-100+ emails daily from clubs, departments, companies, and peers", ACCENT_ORANGE),
    ("🔍", "Critical Emails Lost", "Important deadlines, exam notices, and event invitations get buried under spam and promotions", ACCENT_RED),
    ("⏱️", "Time Wasted", "Hours spent manually sorting through emails instead of focusing on academics", ACCENT_PURPLE),
    ("📋", "No Smart Categorization", "Gmail's default filters are too generic — they can't distinguish academic vs event vs hackathon emails", ACCENT_CYAN),
]

for i, (icon_text, title, desc, color) in enumerate(problems):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(1.8 + (i // 2) * 2.6)
    
    card = add_card(slide, x, y, Inches(5.8), Inches(2.2), border_color=color)
    add_icon_circle(slide, x + Inches(0.3), y + Inches(0.4), Inches(0.8), color, icon_text)
    add_text_box(slide, x + Inches(1.4), y + Inches(0.35), Inches(4.0), Inches(0.5),
                 title, font_size=22, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(1.4), y + Inches(0.9), Inches(4.0), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(-1), Inches(4), Inches(4), Inches(4), ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Our Solution", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "An intelligent system that understands your emails",
             font_size=16, color=ACCENT_GREEN)

# Solution diagram - central hub
hub = add_card(slide, Inches(5.0), Inches(2.5), Inches(3.3), Inches(2.5), 
               color=RGBColor(0x20, 0x20, 0x3A), border_color=ACCENT_BLUE)
add_text_box(slide, Inches(5.0), Inches(2.7), Inches(3.3), Inches(0.5),
             "Smart Email Organizer", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.0), Inches(3.3), Inches(3.3), Inches(1.5),
             "Connect Gmail → Auto-classify → AI Assistant → Actionable insights",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Surrounding feature boxes
features = [
    ("Gmail OAuth 2.0", "Secure read-only access\nto your inbox", Inches(0.5), Inches(2.0), ACCENT_BLUE),
    ("Hybrid ML Engine", "Keyword scoring +\nNaive Bayes classifier", Inches(0.5), Inches(4.5), ACCENT_PURPLE),
    ("AI Chatbot", "Natural language queries\npowered by Llama 3.3", Inches(9.0), Inches(2.0), ACCENT_CYAN),
    ("Smart Calendar", "Visual event timeline\nextracted from emails", Inches(9.0), Inches(4.5), ACCENT_ORANGE),
]

for title, desc, x, y, color in features:
    card = add_card(slide, x, y, Inches(3.5), Inches(1.8), border_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.1), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.1), Inches(0.8),
                 desc, font_size=12, color=LIGHT_GRAY)

# Connecting lines (using thin rectangles)
for y_pos in [Inches(3.2), Inches(5.1)]:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), y_pos, Inches(1.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x60)
    line.line.fill.background()

for y_pos in [Inches(3.2), Inches(5.1)]:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), y_pos, Inches(0.7), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x60)
    line.line.fill.background()

add_slide_number(slide, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 4: KEY FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Key Features", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Everything you need for smarter email management",
             font_size=16, color=ACCENT_BLUE)

features = [
    ("📩", "Gmail Integration", "Secure OAuth 2.0 login\nRead-only inbox access", ACCENT_BLUE),
    ("🤖", "AI Classification", "5-category hybrid ML\nKeyword + Bayes engine", ACCENT_PURPLE),
    ("💬", "Smart Chatbot", "Ask questions in natural language\nPowered by Groq Llama 3.3", ACCENT_CYAN),
    ("📅", "Event Calendar", "Visual timeline of events\nExtracted from email content", ACCENT_ORANGE),
    ("🔍", "Full-Text Search", "Search across subjects,\nsenders, and content", ACCENT_GREEN),
    ("🧠", "Adaptive Learning", "Learns from corrections\nRetrains ML model live", ACCENT_PINK),
]

for i, (emoji, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.7 + row * 2.7)
    
    card = add_card(slide, x, y, Inches(3.8), Inches(2.3), border_color=color)
    add_icon_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.7), color, emoji)
    add_text_box(slide, x + Inches(1.2), y + Inches(0.3), Inches(2.3), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.9),
                 desc, font_size=14, color=LIGHT_GRAY)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: CLASSIFICATION CATEGORIES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(11), Inches(5), Inches(4), Inches(4), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Email Categories", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Five intelligent categories tailored for student life",
             font_size=16, color=ACCENT_PURPLE)

categories = [
    ("📅", "Events", "Campus fests, workshops,\ncultural activities,\nseminars & conferences", CAT_EVENTS),
    ("📚", "Academics", "Assignments, exams,\nscholarships, placements\n& department notices", CAT_ACADEMICS),
    ("💻", "Hackathons", "Coding competitions,\nhackathons, tech\nchallenges & contests", CAT_HACKATHONS),
    ("👤", "Personal", "Friend messages,\nLinkedIn connections,\npersonal notifications", CAT_PERSONAL),
    ("🗑️", "Spam", "Promotions, newsletters,\nautomated notifications\n& marketing emails", CAT_SPAM),
]

for i, (emoji, title, desc, color) in enumerate(categories):
    x = Inches(0.4 + i * 2.5)
    y = Inches(1.8)
    
    card = add_card(slide, x, y, Inches(2.2), Inches(4.5), border_color=color)
    
    # Category icon
    add_icon_circle(slide, x + Inches(0.6), y + Inches(0.4), Inches(1.0), color, emoji)
    
    # Title
    add_text_box(slide, x + Inches(0.1), y + Inches(1.7), Inches(2.0), Inches(0.5),
                 title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_text_box(slide, x + Inches(0.15), y + Inches(2.3), Inches(1.9), Inches(1.8),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: HOW THE CLASSIFIER WORKS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hybrid Classification Engine", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "A multi-layered approach for robust email categorization",
             font_size=16, color=ACCENT_PURPLE)

# Pipeline steps
steps = [
    ("1", "Keyword Scoring\nEngine", "Weighted keyword dictionaries\nper category. Strong=3pts,\nModerate=1pt", ACCENT_BLUE),
    ("2", "Naive Bayes\nClassifier", "NLP text classifier trained\non seed data. Acts as\ntiebreaker for low scores", ACCENT_PURPLE),
    ("3", "Sender-Based\nRules", "Promotional senders → Spam\nAcademic (.edu/.ac.in) →\nAcademics boost", ACCENT_ORANGE),
    ("4", "LinkedIn Special\nHandler", "Distinguishes personal\ninteractions (connections)\nvs notifications (digests)", ACCENT_CYAN),
    ("5", "User Feedback\nLoop", "Corrections retrain Bayes\nmodel & save sender\noverrides for future", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.3 + i * 2.56)
    y = Inches(1.8)
    
    card = add_card(slide, x, y, Inches(2.36), Inches(4.8), border_color=color)
    
    # Step number circle
    num_circle = add_icon_circle(slide, x + Inches(0.78), y + Inches(0.25), Inches(0.7), color, num)
    
    # Title
    add_text_box(slide, x + Inches(0.1), y + Inches(1.2), Inches(2.16), Inches(0.8),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_text_box(slide, x + Inches(0.1), y + Inches(2.3), Inches(2.16), Inches(2.0),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrow between steps
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                        x + Inches(2.36), y + Inches(2.0), Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()

add_slide_number(slide, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: AI CHATBOT
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(-2), Inches(4), Inches(5), Inches(5), ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "AI Chatbot Assistant", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Natural language interaction powered by Groq (Llama 3.3 70B)",
             font_size=16, color=ACCENT_CYAN)

# Left panel - capabilities
cap_card = add_card(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(5.0), border_color=ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.0), Inches(0.5),
             "What you can ask:", font_size=20, color=ACCENT_CYAN, bold=True)

capabilities = [
    "📊  \"Summarize my academic emails\"",
    "📅  \"What deadlines do I have this week?\"",
    "🔍  \"Find emails about hackathons\"",
    "📂  \"Show me my event emails\"",
    "✏️  \"Re-categorize this email as Academic\"",
    "🏠  \"Go back to my dashboard\"",
    "📧  \"How many unread emails do I have?\"",
]

for i, cap in enumerate(capabilities):
    add_text_box(slide, Inches(0.9), Inches(2.7 + i * 0.55), Inches(4.8), Inches(0.5),
                 cap, font_size=14, color=LIGHT_GRAY)

# Right panel - chat mockup
chat_card = add_card(slide, Inches(6.5), Inches(1.8), Inches(6.0), Inches(5.0), 
                     color=RGBColor(0x12, 0x12, 0x22), border_color=RGBColor(0x30, 0x30, 0x50))

# Chat header
add_text_box(slide, Inches(6.8), Inches(2.0), Inches(5.4), Inches(0.4),
             "🤖  AI Assistant", font_size=16, color=ACCENT_CYAN, bold=True)

# Chat bubbles
# User message
user_bubble = add_card(slide, Inches(8.5), Inches(2.6), Inches(3.6), Inches(0.7), 
                       color=RGBColor(0x66, 0x7E, 0xEA))
add_text_box(slide, Inches(8.7), Inches(2.7), Inches(3.2), Inches(0.5),
             "Summarize my academic emails", font_size=12, color=WHITE)

# Bot response
bot_bubble = add_card(slide, Inches(6.8), Inches(3.5), Inches(5.0), Inches(2.8),
                      color=RGBColor(0x20, 0x20, 0x35))
add_text_box(slide, Inches(7.0), Inches(3.6), Inches(4.6), Inches(2.5),
             "📚 Academic Summary:\n\n• Assignment 3 deadline extended to Monday 11:59 PM\n• Deep Learning course registration is open\n• Final exam schedule released — check student portal\n• ISRO VSSC JRF notification posted",
             font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(10), Inches(-1), Inches(5), Inches(5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "Tech Stack", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Modern, production-ready technology choices",
             font_size=16, color=ACCENT_BLUE)

tech_layers = [
    ("Frontend", "React 18 + Vite", "Component-based UI with hot reload,\nReact Router for navigation,\nVanilla CSS with custom properties", ACCENT_BLUE, "⚛️"),
    ("Backend", "Node.js + Express", "RESTful API architecture,\nSession-based authentication,\nModular route handlers", ACCENT_GREEN, "🚀"),
    ("Auth", "Google OAuth 2.0", "Secure Gmail integration,\nRead-only inbox access,\nSession management with cookies", ACCENT_ORANGE, "🔐"),
    ("ML Engine", "natural (Naive Bayes)", "Keyword scoring + NLP classifier,\nTrained on curated seed data,\nAdaptive user feedback loop", ACCENT_PURPLE, "🧠"),
    ("AI / LLM", "Groq (Llama 3.3 70B)", "Ultra-fast inference,\nContext-aware email summarization,\nAction parsing for UI navigation", ACCENT_CYAN, "🤖"),
]

for i, (layer, tech, desc, color, emoji) in enumerate(tech_layers):
    x = Inches(0.5)
    y = Inches(1.7 + i * 1.1)
    
    # Layer label
    label = add_card(slide, x, y, Inches(1.8), Inches(0.9), color=color)
    tf = label.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{emoji} {layer}"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Tech name
    add_text_box(slide, Inches(2.6), y + Inches(0.05), Inches(3.5), Inches(0.4),
                 tech, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.6), y + Inches(0.45), Inches(3.5), Inches(0.5),
                 desc.split('\n')[0], font_size=12, color=LIGHT_GRAY)
    
    # Full description on right
    add_text_box(slide, Inches(7.0), y + Inches(0.1), Inches(5.5), Inches(0.7),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "System Architecture", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "End-to-end data flow from Gmail to user interface",
             font_size=16, color=ACCENT_ORANGE)

# Architecture boxes
arch_items = [
    ("Gmail API", "Fetch emails\nvia OAuth", Inches(0.3), Inches(2.5), ACCENT_BLUE, Inches(2.3)),
    ("Express\nServer", "Routes &\nAPI handlers", Inches(3.1), Inches(2.5), ACCENT_GREEN, Inches(2.3)),
    ("Classifier\nService", "Hybrid ML\nengine", Inches(5.9), Inches(2.5), ACCENT_PURPLE, Inches(2.3)),
    ("Groq LLM", "AI Chatbot\nresponses", Inches(5.9), Inches(4.8), ACCENT_CYAN, Inches(2.3)),
    ("React UI", "Dashboard &\nComponents", Inches(8.7), Inches(2.5), ACCENT_ORANGE, Inches(2.3)),
    ("User", "Interact &\nFeedback", Inches(11.5), Inches(2.5), ACCENT_PINK, Inches(1.5)),
]

for title, desc, x, y, color, w in arch_items:
    card = add_card(slide, x, y, w, Inches(1.8), border_color=color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), Inches(0.6),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), w - Inches(0.2), Inches(0.7),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrows between architecture boxes
arrows_data = [
    (Inches(2.6), Inches(3.3), Inches(0.5)),
    (Inches(5.4), Inches(3.3), Inches(0.5)),
    (Inches(8.2), Inches(3.3), Inches(0.5)),
    (Inches(11.0), Inches(3.3), Inches(0.5)),
]
for ax, ay, aw in arrows_data:
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, aw, Inches(0.2))
    arr.fill.solid()
    arr.fill.fore_color.rgb = RGBColor(0x50, 0x50, 0x70)
    arr.line.fill.background()

# Feedback loop arrow (at bottom)
fb_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.1), Inches(5.5), Inches(9.9), Pt(2))
fb_line.fill.solid()
fb_line.fill.fore_color.rgb = ACCENT_PINK
fb_line.line.fill.background()

add_text_box(slide, Inches(4.5), Inches(5.7), Inches(6), Inches(0.4),
             "← User Feedback Loop: Corrections retrain the classifier & update sender overrides →",
             font_size=12, color=ACCENT_PINK, alignment=PP_ALIGN.CENTER)

# Project structure section
add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
             "📁 Project: client/ (React + Vite)  |  server/ (Node.js + Express)  |  server/services/ (ML + Gmail)  |  server/routes/ (API endpoints)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 10: API ENDPOINTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "API Endpoints", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "RESTful API design with clear separation of concerns",
             font_size=16, color=ACCENT_GREEN)

# Table header
header_card = add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.7), 
                       color=RGBColor(0x25, 0x25, 0x40))
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(1.2), Inches(0.5),
             "Method", font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(2.2), Inches(1.8), Inches(4.5), Inches(0.5),
             "Endpoint", font_size=14, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.5),
             "Description", font_size=14, color=ACCENT_BLUE, bold=True)

api_rows = [
    ("GET", "/auth/google", "Get Google OAuth authentication URL", ACCENT_GREEN),
    ("GET", "/auth/callback", "Handle OAuth callback and create session", ACCENT_GREEN),
    ("GET", "/auth/user", "Get current authenticated user info", ACCENT_GREEN),
    ("POST", "/auth/logout", "Logout and destroy server session", ACCENT_ORANGE),
    ("GET", "/api/emails", "Fetch all emails and classify with ML", ACCENT_BLUE),
    ("GET", "/api/emails/:category", "Fetch emails filtered by category", ACCENT_BLUE),
    ("POST", "/api/emails/:id/recategorize", "Re-categorize email (retrains ML model)", ACCENT_PURPLE),
    ("POST", "/api/chat", "Send message to AI chatbot assistant", ACCENT_CYAN),
]

for i, (method, endpoint, desc, color) in enumerate(api_rows):
    y = Inches(2.5 + i * 0.6)
    row_bg = RGBColor(0x1A, 0x1A, 0x2E) if i % 2 == 0 else RGBColor(0x18, 0x18, 0x28)
    row = add_card(slide, Inches(0.5), y, Inches(12.3), Inches(0.55), color=row_bg)
    
    # Method badge
    method_badge = add_card(slide, Inches(0.7), y + Inches(0.08), Inches(0.9), Inches(0.38), color=color)
    tf = method_badge.text_frame
    p = tf.paragraphs[0]
    p.text = method
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Consolas'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_text_box(slide, Inches(2.2), y + Inches(0.1), Inches(4.5), Inches(0.4),
                 endpoint, font_size=13, color=ACCENT_CYAN, font_name='Consolas')
    add_text_box(slide, Inches(7.0), y + Inches(0.1), Inches(5.5), Inches(0.4),
                 desc, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 11: DEMO & SCREENSHOTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)
add_accent_shape(slide, Inches(10), Inches(4), Inches(5), Inches(5), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "User Interface Highlights", font_size=36, color=WHITE, bold=True, font_name='Segoe UI Semibold')
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
             "Beautiful, responsive dark-themed design with modern aesthetics",
             font_size=16, color=ACCENT_ORANGE)

# UI feature cards
ui_features = [
    ("🌙", "Space-Themed Login", "Animated solar system with orbiting\ncategory planets. Twinkling star field\nbackground with glassmorphism card.", ACCENT_BLUE),
    ("📊", "Bento Grid Dashboard", "Category cards with email counts\nand preview snippets. Hover effects\nwith smooth animations.", ACCENT_GREEN),
    ("📧", "Email Detail View", "Full HTML email rendering with\nattachment support. Inline images\nand download capabilities.", ACCENT_PURPLE),
    ("📅", "Interactive Calendar", "Visual timeline of upcoming events\nextracted from email dates. Click to\nnavigate directly to the email.", ACCENT_ORANGE),
    ("💬", "Floating Chatbot", "Expandable chat widget anchored to\nbottom-right. Real-time AI responses\nwith action buttons.", ACCENT_CYAN),
    ("🔄", "Re-categorize Modal", "One-click email re-classification.\nInstantly updates UI and retrains\nthe ML model in background.", ACCENT_PINK),
]

for i, (emoji, title, desc, color) in enumerate(ui_features):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.2)
    y = Inches(1.7 + row * 2.7)
    
    card = add_card(slide, x, y, Inches(3.9), Inches(2.4), border_color=color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 emoji, font_size=28, color=color)
    add_text_box(slide, x + Inches(0.9), y + Inches(0.25), Inches(2.8), Inches(0.4),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), Inches(3.4), Inches(1.3),
                 desc, font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 12: THANK YOU / CLOSING
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide)

# Large decorative background elements
add_accent_shape(slide, Inches(0), Inches(0), Inches(5), Inches(5), ACCENT_BLUE)
add_accent_shape(slide, Inches(9), Inches(3), Inches(6), Inches(6), ACCENT_PURPLE)
add_accent_shape(slide, Inches(4), Inches(5), Inches(4), Inches(4), ACCENT_CYAN)

# Central card
card = add_card(slide, Inches(2.5), Inches(1.5), Inches(8.3), Inches(4.8),
                color=RGBColor(0x15, 0x15, 0x28), border_color=ACCENT_BLUE)

# Thank you
add_text_box(slide, Inches(2.5), Inches(1.8), Inches(8.3), Inches(1.0),
             "Thank You!", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

add_text_box(slide, Inches(2.5), Inches(2.9), Inches(8.3), Inches(0.5),
             "Smart Email Organizer", font_size=24, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)

# Divider line
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.3), Pt(2))
divider.fill.solid()
divider.fill.fore_color.rgb = ACCENT_BLUE
divider.line.fill.background()

# Summary points
summary_items = [
    "🤖  Hybrid ML engine with 5-category classification",
    "💬  AI Chatbot powered by Groq Llama 3.3",
    "📅  Smart calendar with event extraction",
    "🔄  Adaptive learning from user feedback",
]

for i, item in enumerate(summary_items):
    add_text_box(slide, Inches(3.5), Inches(3.9 + i * 0.45), Inches(6.3), Inches(0.4),
                 item, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom
add_text_box(slide, Inches(2.5), Inches(5.8), Inches(8.3), Inches(0.4),
             "Information Systems Project  •  2026",
             font_size=14, color=SEMI_WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ─── SAVE ───
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Smart_Email_Organizer_Presentation.pptx')
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"     Total slides: {TOTAL_SLIDES}")
